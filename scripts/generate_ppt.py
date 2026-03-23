# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG        = RGBColor(0x0D, 0x0F, 0x1A)
ACCENT    = RGBColor(0x6E, 0xE7, 0xB7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x9C, 0xA3, 0xAF)
YELLOW    = RGBColor(0xFB, 0xBF, 0x24)
RED       = RGBColor(0xF8, 0x71, 0x71)
BLUE      = RGBColor(0x60, 0xA5, 0xFA)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
DARK_CARD = RGBColor(0x1A, 0x1D, 0x2E)
ORANGE    = RGBColor(0xFB, 0x92, 0x3A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def tb(slide, text, left, top, width, height,
       fs=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def accent_bar(slide):
    rect(slide, 0, 0, 13.33, 0.07, ACCENT)

def slide_num(slide, n, total=6):
    tb(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.35,
       fs=10, color=GRAY, align=PP_ALIGN.RIGHT)

def tag(slide, text, color=ACCENT):
    rect(slide, 0.5, 0.22, 2.6, 0.38, color)
    tb(slide, text, 0.56, 0.24, 2.5, 0.33, fs=11, bold=True,
       color=RGBColor(0x0D, 0x0F, 0x1A))

def arrow(slide, x, y):
    tb(slide, "->", x, y, 0.45, 0.45, fs=22, bold=True,
       color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 0 — Title / Intro
# ══════════════════════════════════════════════════════════════════
s0 = blank_slide(prs)
set_bg(s0)
rect(s0, 0, 0, 13.33, 0.07, ACCENT)
rect(s0, 0, 7.43, 13.33, 0.07, ACCENT)

# Subtle background grid lines
for gx in [3.33, 6.66, 9.99]:
    rect(s0, gx, 0, 0.01, 7.5, RGBColor(0x1A, 0x1D, 0x2E))

# Tag
rect(s0, 0.55, 0.55, 3.4, 0.42, RGBColor(0x1A, 0x1D, 0x2E))
tb(s0, "FLEXIBLE STREET PLATFORM  —  ML SCORING",
   0.65, 0.58, 3.25, 0.35, fs=10, bold=False, color=ACCENT)

# Main headline
tb(s0, "Using Machine Learning\nto Score Street Openings",
   0.55, 1.35, 11.5, 2.2, fs=52, bold=True, color=WHITE)

# Subtitle
tb(s0,
   "We apply XGBoost to evaluate every anchor point on the Philadelphia street network\n"
   "and predict whether its surrounding street is suitable for temporary public activation.",
   0.55, 3.85, 10.5, 1.0, fs=18, color=GRAY)

# Three pill tags
pills = [
    ("7 Scenario Types",    BLUE),
    ("XGBoost + SHAP",      ACCENT),
    ("Open Data Sources",   YELLOW),
]
for pi, (label, col) in enumerate(pills):
    px = 0.55 + pi * 3.6
    rect(s0, px, 5.1, 3.3, 0.55, RGBColor(0x1A, 0x1D, 0x2E))
    rect(s0, px, 5.1, 0.07, 0.55, col)
    tb(s0, label, px + 0.22, 5.18, 2.95, 0.38, fs=14, bold=True, color=col)

# Bottom line
tb(s0, "Philadelphia, PA  |  Urban Planning  |  Spring 2026",
   0.55, 6.9, 8.0, 0.38, fs=11, color=GRAY, italic=True)

slide_num(s0, 1)

notes(s0, """Good afternoon everyone. Today we're presenting the ML scoring component of our Flexible Street Platform.

The core question is: given a specific anchor point on the Philadelphia street network — a school, a restaurant, a transit hub — can we use machine learning to predict whether the street around it is suitable for temporary public activation?

We use XGBoost as our scoring model, applied across seven different urban scenario types, drawing entirely from Philadelphia's existing open data. The output is an FSI score from 0 to 100, a three-tier recommendation, and a SHAP-based explanation of the key contributing factors.

Let's walk through how it works.""")

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Why XGBoost?
# ══════════════════════════════════════════════════════════════════
s1 = blank_slide(prs)
set_bg(s1); accent_bar(s1); slide_num(s1, 2)
tag(s1, "MODEL SELECTION", BLUE)

tb(s1, "Why XGBoost?", 0.55, 0.75, 12, 0.65, fs=36, bold=True)
tb(s1, "More accurate than Random Forest, still interpretable, ideal for small structured datasets",
   0.55, 1.45, 12, 0.4, fs=15, color=GRAY)

# Left card — Random Forest
rect(s1, 0.35, 2.0, 5.9, 4.85, DARK_CARD)
tb(s1, "Random Forest  (baseline)", 0.55, 2.1, 5.5, 0.42,
   fs=14, bold=True, color=GRAY)
for i, line in enumerate([
    "500 trees grow independently in parallel",
    "Each tree sees a random data subset",
    "Final prediction = average of all trees",
    "Supports SHAP interpretability",
]):
    tb(s1, f"  {line}", 0.6, 2.68 + i*0.55, 5.5, 0.48, fs=13, color=GRAY)

# Divider arrow
tb(s1, "vs", 6.3, 4.1, 0.7, 0.55, fs=20, bold=True,
   color=ACCENT, align=PP_ALIGN.CENTER)

# Right card — XGBoost
rect(s1, 6.55, 2.0, 6.4, 4.85, DARK_CARD)
rect(s1, 6.55, 2.0, 6.4, 0.07, ACCENT)
tb(s1, "XGBoost  (recommended)", 6.75, 2.1, 6.0, 0.42,
   fs=14, bold=True, color=ACCENT)
for i, line in enumerate([
    "Trees built sequentially, not in parallel",
    "Each new tree corrects the previous tree's errors",
    "Gradient descent minimizes prediction loss",
    "Higher accuracy on same data volume",
    "Supports SHAP interpretability",
]):
    tb(s1, f"  {line}", 6.75, 2.68 + i*0.55, 6.1, 0.48, fs=13, color=ACCENT)

# Bottom two notes
rect(s1, 0.35, 5.08, 6.05, 1.6, RGBColor(0x10, 0x15, 0x25))
tb(s1, "Core upgrade", 0.55, 5.18, 5.7, 0.35, fs=12, bold=True)
tb(s1, "RF: all trees independent, errors remain\n"
       "XGBoost: each tree focuses on what the\n"
       "previous tree got wrong -> error shrinks\n"
       "with every iteration",
   0.55, 5.55, 5.7, 1.05, fs=11, color=GRAY)

rect(s1, 6.55, 5.08, 6.4, 1.6, RGBColor(0x10, 0x15, 0x25))
tb(s1, "Why not Deep Learning?", 6.75, 5.18, 6.0, 0.35, fs=12, bold=True)
tb(s1, "Philadelphia Play Streets labels are limited\n"
       "(hundreds of records, not millions).\n"
       "Deep learning needs massive data.\n"
       "XGBoost dominates small structured datasets.",
   6.75, 5.55, 6.1, 1.05, fs=11, color=GRAY)

notes(s1, """We chose XGBoost as our scoring model instead of Random Forest.

Both are tree-based ensemble methods, but the logic is different. Random Forest grows 500 trees independently in parallel and averages their results. XGBoost builds trees sequentially — each new tree is specifically designed to correct the errors made by the previous one, using gradient descent to minimize the prediction loss step by step.

Both models support SHAP interpretability, so planners can still understand why a street gets a certain score.

We did not choose deep learning because Philadelphia's Play Streets historical labels are limited to a few hundred records. Deep learning requires massive datasets to generalize well and tends to overfit on small ones. XGBoost is the industry-standard choice for small structured tabular data.""")

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Anchor Differentiation
# ══════════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
set_bg(s2); accent_bar(s2); slide_num(s2, 3)
tag(s2, "ANCHOR DIFFERENTIATION", YELLOW)

tb(s2, "Different Anchors, Different Scoring Logic",
   0.55, 0.75, 12, 0.65, fs=34, bold=True)
tb(s2, "7 scenario types — each trained as an independent XGBoost model with its own feature weights",
   0.55, 1.45, 12, 0.4, fs=15, color=GRAY)

scenarios = [
    ("School\nDismissal",  BLUE,
     "Key: time_match\n+ pedestrian_safety",
     "Public: school bus dominant\n-> step timing predictable\nPrivate: car pickup heavy\n-> parking conflict weight up",
     "school_type=0/1\nbranches weight"),
    ("Weekend\nMarket",    YELLOW,
     "Key: foot_traffic\n+ is_food_desert",
     "Weekend crowd density +\nfood desert community flag\nDetour capacity of\nsurrounding streets",
     "food_desert as\nequity modifier"),
    ("Community\nEvent",   PURPLE,
     "Key: schedule\n+ detour_capacity",
     "Does event have a fixed\npredictable schedule?\nCan nearby roads absorb\ndiverted traffic?",
     "Trained separately\nfrom Weekend Market"),
    ("Dining\nActivation", RED,
     "Key: is_evening\n+ weather_score",
     "Evening slot + temp/rain\nscore + sports event flag\nExclude rush hours and\ndelivery-dependent windows",
     "season_score as\nmultiplication modifier"),
    ("Play\nStreet",       ACCENT,
     "Key: child_density\n+ traffic_volume",
     "0-14 yr population density\n+ daily vehicle count\nHistoric Play Streets\napproval = strong signal",
     "has_play_street_history\nhighest weight"),
    ("Event\nSpillover",   ORANGE,
     "Key: crowd_size\n+ venue_distance",
     "Venue capacity x attendance\n+ distance to venue\nArterial conflict check\nbefore model entry",
     "arterial_conflict\nhard veto"),
    ("Transit\nZone",      RGBColor(0x38,0xBD,0xF8),
     "Key: transit_freq\n+ pedestrian_overflow",
     "30th St Station example:\nAmtrak + SEPTA transfer density\nRush hour = hard veto\nbefore model entry",
     "bus_route_dependency\nas hard constraint"),
]

cw = 12.6 / 7
for i, (name, col, feat, logic, note_txt) in enumerate(scenarios):
    x = 0.35 + i * cw
    # header
    rect(s2, x, 1.92, cw - 0.07, 0.65, col)
    tb(s2, name, x + 0.08, 1.96, cw - 0.17, 0.58,
       fs=9, bold=True, color=RGBColor(0x0D, 0x0F, 0x1A))
    # body
    rect(s2, x, 2.57, cw - 0.07, 4.53, DARK_CARD)
    rect(s2, x, 2.57, 0.05, 4.53, col)
    tb(s2, feat,     x+0.12, 2.64, cw-0.22, 0.62, fs=8.5, bold=True, color=col)
    tb(s2, logic,    x+0.12, 3.32, cw-0.22, 1.65, fs=8,   color=WHITE)
    tb(s2, note_txt, x+0.12, 5.05, cw-0.22, 0.85, fs=7.5, color=GRAY, italic=True)

rect(s2, 0.35, 6.28, 12.6, 0.82, RGBColor(0x12, 0x15, 0x22))
tb(s2, "Shared mechanism:",
   0.55, 6.37, 1.8, 0.55, fs=11, bold=True)
tb(s2, "All 7 scenarios share the same feature schema (temporal / activity / safety / community). "
       "Differences lie in per-scenario feature weights. "
       "Hard constraints (emergency routes, peak hours) are filtered before model entry.",
   2.4, 6.37, 10.4, 0.58, fs=11, color=GRAY)

notes(s2, """This slide explains how we handle the seven different anchor types differently.

The core idea is that all seven scenarios share the same feature schema, but each scenario has an independent XGBoost model with different feature weights.

Take School Dismissal as an example. We split it into public and private schools. Public schools are mainly served by school buses — dismissal timing is standardized and predictable, so pedestrian safety weight is higher. Private schools have heavy car pickup, so parking conflict features carry more weight.

Transit Zone uses 30th Street Station as a reference — the Amtrak and SEPTA transfer density is the key feature, and rush hour is a hard veto before the model even runs.

The bottom row summarizes the shared mechanism: same feature schema, differentiated weights, hard constraints applied before model entry.""")

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — XGBoost Training Logic Chain
# ══════════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
set_bg(s3); accent_bar(s3); slide_num(s3, 4)
tag(s3, "XGBOOST TRAINING LOGIC", ACCENT)

tb(s3, "How XGBoost Learns — The Boosting Chain",
   0.55, 0.75, 12, 0.65, fs=34, bold=True)
tb(s3, "Each tree is built to fix what the previous tree got wrong — error shrinks at every step",
   0.55, 1.45, 12, 0.4, fs=15, color=GRAY)

# ── Training data row ──
rect(s3, 0.35, 1.95, 12.6, 0.06, RGBColor(0x22, 0x25, 0x38))

rect(s3, 0.35, 2.08, 5.9, 0.78, RGBColor(0x14, 0x3A, 0x1E))
tb(s3, "Positive samples  (label = 1  ->  should open)",
   0.55, 2.15, 5.6, 0.35, fs=12, bold=True, color=ACCENT)
tb(s3, "Philadelphia Play Streets Program — officially approved street closures",
   0.55, 2.5, 5.6, 0.3, fs=11, color=GRAY)

rect(s3, 6.6, 2.08, 6.38, 0.78, RGBColor(0x4A, 0x1A, 0x1A))
tb(s3, "Negative samples  (label = 0  ->  should NOT open)",
   6.78, 2.15, 6.05, 0.35, fs=12, bold=True, color=RED)
tb(s3, "Vision Zero high-accident corridors  +  peak-hour arterials",
   6.78, 2.5, 6.0, 0.3, fs=11, color=GRAY)

# ── Boosting chain ──
chain_items = [
    ("Tree 1",     "Learns from\nraw features",       "First pass prediction\nLarge residual errors remain",   BLUE),
    ("Residuals",  "Prediction errors\nfrom Tree 1",   "Which anchors did\nTree 1 get wrong?",                 GRAY),
    ("Tree 2",     "Trained on\nresiduals of Tree 1",  "Focuses on hard cases\nError shrinks",                 RGBColor(0x34,0x8C,0xF7)),
    ("Residuals",  "Remaining errors\nafter Tree 2",   "Even smaller residuals\npassed forward",               GRAY),
    ("Tree N",     "...500 trees\nlater...",            "Cumulative correction\nHigh accuracy achieved",        ACCENT),
]

box_w, box_h = 2.15, 1.55
gap = 0.18
total_w = len(chain_items) * box_w + (len(chain_items)-1) * (gap + 0.38)
start_x = (13.33 - total_w) / 2

for i, (title, sub, caption, col) in enumerate(chain_items):
    x = start_x + i * (box_w + gap + 0.38)
    is_residual = title == "Residuals"
    bg = RGBColor(0x22,0x25,0x38) if is_residual else DARK_CARD
    rect(s3, x, 3.05, box_w, box_h, bg)
    rect(s3, x, 3.05, box_w, 0.07, col)
    tb(s3, title,   x+0.12, 3.14, box_w-0.22, 0.42, fs=13, bold=True, color=col)
    tb(s3, sub,     x+0.12, 3.58, box_w-0.22, 0.55, fs=10, color=WHITE)
    tb(s3, caption, x+0.12, 4.18, box_w-0.22, 0.72, fs=9,  color=GRAY, italic=True)
    if i < len(chain_items) - 1:
        ax = x + box_w + 0.04
        tb(s3, "->", ax, 3.6, 0.35, 0.42, fs=20, bold=True,
           color=ACCENT, align=PP_ALIGN.CENTER)

# ── Loss curve description ──
rect(s3, 0.35, 4.78, 6.1, 2.42, DARK_CARD)
rect(s3, 0.35, 4.78, 6.1, 0.07, YELLOW)
tb(s3, "Loss Curve (conceptual)", 0.55, 4.88, 5.7, 0.38, fs=13, bold=True, color=YELLOW)
tb(s3, "Prediction Error", 0.45, 5.32, 1.3, 1.6, fs=10, color=GRAY,
   align=PP_ALIGN.CENTER, italic=True)

# Draw simple stepdown bars to simulate loss curve
bar_heights = [1.35, 1.0, 0.72, 0.52, 0.38, 0.28, 0.22]
bar_cols = [RGBColor(0xF8,0x71,0x71), RGBColor(0xF9,0x90,0x55),
            RGBColor(0xFB,0xBF,0x24), RGBColor(0xA3,0xE6,0x35),
            RGBColor(0x6E,0xE7,0xB7), RGBColor(0x4A,0xDE,0x80),
            RGBColor(0x22,0xC5,0x5E)]
for bi, (bh, bc) in enumerate(zip(bar_heights, bar_cols)):
    bx = 1.85 + bi * 0.6
    by = 5.28 + (1.35 - bh)
    rect(s3, bx, by, 0.48, bh, bc)

tb(s3, "Tree 1   Tree 2   Tree 3   ...   Tree N",
   1.72, 6.72, 4.4, 0.35, fs=9, color=GRAY, italic=True)

# ── Feature importance ──
rect(s3, 6.6, 4.78, 6.38, 2.42, DARK_CARD)
rect(s3, 6.6, 4.78, 6.38, 0.07, PURPLE)
tb(s3, "Feature Importance (School Dismissal example)",
   6.78, 4.88, 6.1, 0.38, fs=13, bold=True, color=PURPLE)

feat_imp = [
    ("time_match_score",        0.88, ACCENT),
    ("pedestrian_density",      0.72, ACCENT),
    ("child_density",           0.61, BLUE),
    ("accident_count",          0.45, YELLOW),
    ("has_bus_route",           0.38, YELLOW),
    ("car_pickup_pressure",     0.25, RED),
]
bar_total = 5.5
for fi, (fname, imp, col) in enumerate(feat_imp):
    fy = 5.38 + fi * 0.33
    tb(s3, fname, 6.78, fy, 2.1, 0.28, fs=9, color=GRAY)
    rect(s3, 8.95, fy+0.04, bar_total * imp, 0.2, col)
    tb(s3, f"{int(imp*100)}%", 8.98 + bar_total * imp, fy,
       0.55, 0.28, fs=9, color=WHITE)

notes(s3, """This slide digs into how XGBoost actually learns.

We start with two sets of training data. Positive samples come from the Philadelphia Play Streets Program — streets that have been officially approved for temporary closure. These are labeled 1. Negative samples come from Vision Zero high-accident corridors and peak-hour arterials, labeled 0.

The boosting chain works like this. Tree 1 learns from the raw feature vectors and makes a first-pass prediction. It will be wrong on many samples — those errors are called residuals. Tree 2 is then trained specifically on those residuals, focusing on the cases Tree 1 got wrong. This produces a smaller set of residuals. The process repeats for up to 500 trees, with each tree correcting the previous one's mistakes. By the end, the cumulative prediction is highly accurate.

The loss curve on the bottom left shows this visually — prediction error drops steeply in the first few trees and gradually flattens as the model converges.

The bottom right shows feature importance for the School Dismissal scenario. Time match score has the highest importance, followed by pedestrian density and child density. This tells us which features XGBoost relied on most heavily during training.""")

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — Full Model Pipeline
# ══════════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
set_bg(s4); accent_bar(s4); slide_num(s4, 5)
tag(s4, "MODEL PIPELINE", ACCENT)

tb(s4, "End-to-End Model Pipeline",
   0.55, 0.75, 12, 0.65, fs=34, bold=True)
tb(s4, "From raw open data to a street-opening decision — four steps",
   0.55, 1.45, 10, 0.4, fs=15, color=GRAY)

steps = [
    (BLUE,   "Step 1\nData Ingestion",
     "OpenDataPhilly\nSEPTA GTFS\nACS Census\nOpenStreetMap\nVision Zero\nOpenWeather API",
     "6 public databases\nSpatial join to\neach anchor point"),
    (YELLOW, "Step 2\nFeature Engineering",
     "Temporal\n  time_match_score\n  is_rush_hour\n\nActivity\n  pedestrian_density\n  restaurant_density\n\nSafety\n  accident_count  lts_score\n\nCommunity\n  child_density\n  is_food_desert",
     "Each anchor becomes\na numeric feature vector\nHard constraints filtered\nbefore model entry"),
    (ACCENT, "Step 3\nXGBoost Training",
     "Positive labels\n  Play Streets approvals\n\nNegative labels\n  Vision Zero + peak\n  hour arterials\n\n7 independent models\n  one per scenario\n\nGradient boosting\n  ~500 trees each",
     "Error corrected\niteration by iteration\nOutputs FSI score 0-100"),
    (RED,    "Step 4\nResult Output",
     "FSI Score  0 - 100\n\nRecommendation\n  Open          70-100\n  Conditional   40-69\n  Not advised    0-39\n\nSHAP factor breakdown",
     "Written to Supabase\nDisplayed in frontend\nAnchorDetailPanel"),
]

for i, (col, title, body, caption) in enumerate(steps):
    x = 0.35 + i * 3.25
    rect(s4, x, 1.92, 3.1, 1.0, col)
    tb(s4, title, x+0.15, 1.97, 2.85, 0.88, fs=17, bold=True,
       color=RGBColor(0x0D, 0x0F, 0x1A), align=PP_ALIGN.CENTER)
    rect(s4, x, 2.95, 3.1, 3.55, DARK_CARD)
    rect(s4, x, 2.95, 0.06, 3.55, col)
    tb(s4, body,    x+0.18, 3.05, 2.85, 3.35, fs=11, color=WHITE)
    rect(s4, x, 6.53, 3.1, 0.72, RGBColor(0x10, 0x12, 0x1E))
    tb(s4, caption, x+0.15, 6.58, 2.85, 0.62, fs=10, color=col, italic=True)
    if i < 3:
        tb(s4, "->", x+3.15, 4.0, 0.32, 0.45, fs=22, bold=True,
           color=ACCENT, align=PP_ALIGN.CENTER)

tb(s4, "* Hard constraints (emergency routes / rush-hour arterials) "
       "are filtered at Step 2 before entering the model",
   0.55, 7.2, 12.5, 0.28, fs=10, color=GRAY, italic=True)

notes(s4, """This slide shows the complete four-step pipeline.

Step 1: Data ingestion. We pull from six public databases and use spatial joins to attach the relevant data to each anchor point on the map.

Step 2: Feature engineering. We compute four categories of features for each anchor: temporal, activity, safety, and community. Hard constraints like emergency routes and rush-hour arterials are filtered out at this stage before the model runs.

Step 3: XGBoost training. Positive labels come from Play Streets approvals, negative labels from Vision Zero and peak-hour arterials. Each of the seven scenarios trains its own independent XGBoost model with roughly 500 trees using gradient boosting.

Step 4: Output. The FSI score is written to Supabase and displayed in the frontend. Each anchor shows a score from 0 to 100, a three-tier recommendation, and a SHAP breakdown of the main contributing factors.""")

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Frontend Output
# ══════════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
set_bg(s5); accent_bar(s5); slide_num(s5, 6)
tag(s5, "FRONTEND OUTPUT", PURPLE)

tb(s5, "Frontend Output", 0.55, 0.75, 12, 0.65, fs=36, bold=True)
tb(s5, "Click any anchor on the map to see its FSI score, recommendation, and SHAP explanation",
   0.55, 1.45, 12, 0.4, fs=15, color=GRAY)

# ── Mock AnchorDetailPanel ──
rect(s5, 0.35, 2.0, 4.85, 5.22, RGBColor(0x13, 0x14, 0x1C))
rect(s5, 0.35, 2.0, 4.85, 0.06, ACCENT)

tb(s5, "Clark Elementary School", 0.55, 2.12, 4.5, 0.46, fs=17, bold=True)
tb(s5, "School Dismissal   |   amenity: school",
   0.55, 2.6, 4.5, 0.35, fs=11, color=GRAY)

rect(s5, 0.45, 3.05, 4.65, 1.08, DARK_CARD)
tb(s5, "82", 0.6, 3.07, 0.95, 1.0, fs=44, bold=True, color=ACCENT)
tb(s5, "/ 100", 1.52, 3.44, 0.8, 0.38, fs=14, color=GRAY)
rect(s5, 2.5, 3.22, 2.45, 0.5, RGBColor(0x14, 0x3A, 0x1E))
tb(s5, "Recommend Open  v", 2.58, 3.3, 2.3, 0.35, fs=13, bold=True, color=ACCENT)
rect(s5, 0.45, 4.19, 4.65, 0.2, RGBColor(0x2A, 0x2D, 0x3E))
rect(s5, 0.45, 4.19, 4.65*0.82, 0.2, ACCENT)

tb(s5, "SHAP Factor Breakdown", 0.55, 4.52, 4.4, 0.35, fs=12, bold=True)
shap = [
    ("+35", "Dismissal time match",    ACCENT),
    ("+28", "Low traffic conflict",    ACCENT),
    ("+19", "High child density",      ACCENT),
    ("-12", "SEPTA bus route nearby",  RED),
    (" -8", "Moderate parking demand", RED),
]
for j, (sc, lb, col) in enumerate(shap):
    rect(s5, 0.45, 4.92 + j*0.38, 4.65, 0.34, RGBColor(0x15, 0x18, 0x28))
    tb(s5, sc, 0.52, 4.95+j*0.38, 0.6, 0.28, fs=11, bold=True, color=col)
    tb(s5, lb, 1.18, 4.95+j*0.38, 3.8, 0.28, fs=11, color=WHITE)

# ── Three tiers ──
tiers = [
    ("70 – 100", "Recommend\nOpen",        "Conditions met\nApply directly",         ACCENT, RGBColor(0x14,0x3A,0x1E)),
    ("40 – 69",  "Conditional\nOpen",      "Requires traffic\ncontrol measures",      YELLOW, RGBColor(0x3A,0x2E,0x0A)),
    ("0 – 39",   "Not\nAdvisable",         "Risk too high\ndo not open",              RED,    RGBColor(0x4A,0x1A,0x1A)),
]
for k, (rng, label, desc, col, bg) in enumerate(tiers):
    x = 5.55 + k * 2.62
    rect(s5, x, 2.0, 2.5, 2.7, bg)
    rect(s5, x, 2.0, 2.5, 0.06, col)
    tb(s5, rng,   x+0.15, 2.12, 2.25, 0.45, fs=15, bold=True, color=col)
    tb(s5, label, x+0.15, 2.6,  2.25, 0.62, fs=20, bold=True, color=WHITE)
    tb(s5, desc,  x+0.15, 3.28, 2.25, 0.85, fs=12, color=GRAY)

# ── Why SHAP ──
rect(s5, 5.55, 4.85, 7.5, 2.38, DARK_CARD)
rect(s5, 5.55, 4.85, 7.5, 0.07, PURPLE)
tb(s5, "Why SHAP Explainability Matters",
   5.72, 4.96, 7.15, 0.4, fs=14, bold=True, color=PURPLE)
tb(s5,
   "Planners don't just need a number — they need to justify the decision "
   "to government approval bodies and community residents.\n\n"
   "SHAP (SHapley Additive exPlanations) quantifies each feature's contribution "
   "to the final score as a positive or negative value, making the model's reasoning "
   "transparent, communicable, and defensible.",
   5.72, 5.44, 7.22, 1.72, fs=12, color=GRAY)

notes(s5, """The final slide shows what planners actually see in the frontend.

When a planner clicks any anchor on the map, a detail panel appears on the right side of the screen showing three sections.

The first section is the FSI score and recommendation. In this example, Clark Elementary School scores 82 out of 100, placing it in the Recommend Open tier.

The second section is the SHAP factor breakdown. This explains how the 82 was reached: dismissal time match contributes plus 35, low traffic conflict adds 28, high child density adds 19, but a nearby SEPTA bus route subtracts 12, and moderate parking demand subtracts 8.

The third section explains why SHAP matters. Planners need to be able to explain their decisions to city government and community members. SHAP turns the model's internal reasoning into plain positive and negative contributions for each feature, making the decision transparent and defensible.

The three tiers at the top right define the thresholds: 70 to 100 is recommend open, 40 to 69 requires traffic control measures, and 0 to 39 means the street should not be opened.""")

# ── Save ──
out = r"e:\Upenn_course\Spring\flexible-street-platform\Flexible_Street_ML_Scoring_v3.pptx"
prs.save(out)
print(f"Saved -> {out}")
